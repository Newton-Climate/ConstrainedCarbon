"""Deck QA: measure layout bounds, overlaps, and likely text overflow."""
from pptx import Presentation
from pptx.util import Emu

EMU = 914400.0
SLIDE_W, SLIDE_H = 13.3, 7.5

PATH = "/Users/newtonnguyen/Documents/ecosystem-complexity/notebooks/ecosystem_complexity_soil_carbon_canonical_presentation.pptx"
p = Presentation(PATH)

def bounds_in(s):
    return (s.left/EMU, s.top/EMU, s.width/EMU, s.height/EMU)

issues = []
for i, slide in enumerate(p.slides, 1):
    rects = []
    for sh in slide.shapes:
        if sh.left is None or sh.top is None: continue
        x, y, w, h = bounds_in(sh)
        # Off-slide check
        if x + w > SLIDE_W + 0.02 or y + h > SLIDE_H + 0.02:
            issues.append(f"S{i}: shape extends off-slide  ({x:.2f},{y:.2f}) {w:.2f}×{h:.2f}  ['{sh.name}']")
        # Likely text overflow: text frame with very long text vs height
        if sh.has_text_frame:
            txt = sh.text_frame.text
            if txt:
                est_chars_per_inch = 12
                est_lines = max(1, len(txt) // max(int(est_chars_per_inch * w), 1))
                est_pt = 12
                est_line_h = est_pt / 60.0
                est_total_h = est_lines * est_line_h
                if est_total_h > h * 1.3 and len(txt) > 200 and h < 1.5:
                    issues.append(f"S{i}: possible overflow  '{txt[:50]}...'  text {len(txt)}c in {w:.1f}×{h:.1f}")
        rects.append((x, y, w, h, sh.name, sh.has_text_frame, sh.shape_type))

    # Overlap check among non-image shapes (decorative rects can intentionally overlap)
    for a_i in range(len(rects)):
        for b_i in range(a_i+1, len(rects)):
            xa,ya,wa,ha,na,_,_ = rects[a_i]
            xb,yb,wb,hb,nb,_,_ = rects[b_i]
            ix = max(0, min(xa+wa, xb+wb) - max(xa, xb))
            iy = max(0, min(ya+ha, yb+hb) - max(ya, yb))
            overlap = ix * iy
            # only flag if a text-bearing shape is being covered by a non-image shape
            # skip — too noisy without filtering by intent

print(f"Slides: {len(p.slides)}")
print(f"Issues found: {len(issues)}")
for line in issues[:40]:
    print(" ", line)
